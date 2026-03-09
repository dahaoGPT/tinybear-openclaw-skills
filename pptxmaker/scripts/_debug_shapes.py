"""诊断脚本：详细列出每页每个形状的类型、位置、文本"""
from pptx import Presentation
from pptx.util import Inches

TEMPLATE = r'e:\2-AI\6-workspaces\managertools\claude-code\pptxmaker\templates\company-2026.pptx'
prs = Presentation(TEMPLATE)

for slide_idx, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else '?'
    print(f"\n{'='*60}")
    print(f"第{slide_idx+1}页 [布局: {layout_name}]")
    print(f"{'='*60}")

    # 先列出占位符
    phs = list(slide.placeholders)
    if phs:
        print(f"  📌 占位符 ({len(phs)}个):")
        for ph in phs:
            txt = ph.text_frame.text[:80].replace('\n', '|') if ph.has_text_frame else '(无文本框)'
            print(f"    idx={ph.placeholder_format.idx} name='{ph.name}' type={ph.placeholder_format.type}")
            print(f"      文本: '{txt}'")
    else:
        print("  📌 占位符: 无")

    # 列出所有非占位符的形状
    non_ph_shapes = []
    ph_names = {ph.name for ph in phs}
    for shape in slide.shapes:
        if shape.name not in ph_names:
            non_ph_shapes.append(shape)

    if non_ph_shapes:
        print(f"  📦 其他形状 ({len(non_ph_shapes)}个):")
        for shape in non_ph_shapes:
            txt = ''
            if shape.has_text_frame:
                txt = shape.text_frame.text[:80].replace('\n', '|')
            shape_type = type(shape).__name__
            print(f"    name='{shape.name}' type={shape_type}")
            if txt:
                print(f"      文本: '{txt}'")
    
    # 只看前 5 页 + 分隔页样本
    if slide_idx >= 5:
        print("\n...(省略后续页面，结构类似)...")
        break

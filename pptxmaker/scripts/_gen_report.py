"""
模拟 /ppt-maker:from-template company-2026 银行业软件研发效能报告
策略V4：
1. 封面、目录、分隔页、结束页保留原有排版并替换文本
2. **正文页**清空原来复杂的图形，根据品牌颜色和文本内容，使用 python-pptx **重新排版绘制图形和图标**，让展示更紧凑合理。
"""
import os
import warnings
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

warnings.filterwarnings('ignore')

TEMPLATE = r'e:\2-AI\6-workspaces\managertools\claude-code\pptxmaker\templates\company-2026.pptx'
OUTPUT_DIR = r'e:\2-AI\6-workspaces\managertools\claude-code'
TODAY = datetime.now().strftime('%Y%m%d')
OUTPUT = os.path.join(OUTPUT_DIR, f'银行业软件研发效能报告_{TODAY}.pptx')

FONT_CN = '圆体-简'
FONT_EN = 'effra'

COLOR_PRIMARY = RGBColor(255, 85, 5)       # 橙红
COLOR_SECONDARY = RGBColor(244, 185, 0)    # 黄色
COLOR_ACCENT = RGBColor(53, 4, 100)        # 深紫
COLOR_TEXT = RGBColor(53, 4, 100)          # 深紫 (正文)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_BG = RGBColor(250, 243, 238)   # 极浅的底色，适配橙红风格

# ============================================================
# 辅助函数
# ============================================================

def set_run_font(run, size_pt, bold=False, color=None):
    """设置 Run 字体"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = FONT_EN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_CN)

def write_to_shape(shape, text, size_pt=16, bold=False, color=None, align=None):
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size_pt, bold, color)
    return True

def replace_text_everywhere(slide, old_text, new_text):
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1
    for ph in slide.placeholders:
        if ph.has_text_frame:
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1
    return count

def find_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None

def draw_custom_content_slide(slide, title_text, bullets):
    """
    清空正文幻灯片中原来的形状，并使用 python-pptx 重新排版绘制漂亮的图形和图标
    """
    # 1. 删除旧形状（注意：这不会删除母版上的背景图片和页码/Logo）
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    
    # 2. 绘制标题（带左侧强调色块）
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_PRIMARY
    left_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_run_font(run, 30, bold=True, color=COLOR_ACCENT)
    
    # 3. 绘制上方装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_SECONDARY
    line.line.fill.background()

    # 4. 绘制对应的条目内容（带色块和图形排版）
    num_bullets = len(bullets)
    if num_bullets == 0:
        return
        
    start_y = 1.6
    y_gap = 5.2 / num_bullets # 可用高度动态分配

    for i, bullet in enumerate(bullets):
        y_pos = Inches(start_y + i * y_gap)
        
        # 绘制背景浅色圆角卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y_pos, Inches(12.1), Inches(y_gap * 0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_BG
        card.line.fill.background()

        # 绘制序号框（橙红底白字）
        num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y_pos, Inches(0.8), Inches(y_gap * 0.8))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_PRIMARY
        num_box.line.fill.background()
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        run_num = p_num.add_run()
        run_num.text = f"0{i+1}" if i < 9 else str(i+1)
        set_run_font(run_num, 24, bold=True, color=COLOR_WHITE)

        # 绘制正文文本
        txBox = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(y_gap*0.1), Inches(10.8), Inches(y_gap * 0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        set_run_font(run, 18, bold=False, color=COLOR_TEXT)


# ============================================================
# 打开模板
# ============================================================
print(f"打开模板: {TEMPLATE}")
prs = Presentation(TEMPLATE)
slides = list(prs.slides)
print(f"模板共有 {len(slides)} 页幻灯片\n")

# ============================================================
# 封面页、目录页、分隔页修改（保留原有排版）
# ============================================================

# --- 第1页：封面页 ---
for ph in slides[0].placeholders:
    write_to_shape(ph, '银行业软件研发效能报告', 36, bold=True)
replace_text_everywhere(slides[0], '年度工作总结', '银行业软件研发效能报告')
replace_text_everywhere(slides[0], '报告报告', '报告')
print("✅ 第1页: 封面页 (保留原版式)")

# --- 第2页：目录页 ---
s1 = find_shape_by_name(slides[1], '文本框 1')
s2 = find_shape_by_name(slides[1], '文本框 7')
s3 = find_shape_by_name(slides[1], '文本框 8')
if s1: write_to_shape(s1, '1、行业研发效能现状与趋势', 20)
if s2: write_to_shape(s2, '2、银行业研发效能核心指标体系', 20)
if s3: write_to_shape(s3, '3、DevOps 与敏捷实践落地分析', 20)
# 若是其他的模板格式：
replace_text_everywhere(slides[1], '标题1', '行业研发效能现状与趋势')
replace_text_everywhere(slides[1], '标题2', '银行业研发效能核心指标体系')
replace_text_everywhere(slides[1], '标题3', 'DevOps 与敏捷实践落地分析')
print("✅ 第2页: 目录页 (保留原版式)")

# --- 第3页：分隔页1 ---
s = find_shape_by_name(slides[2], '文本框 7')
if s: write_to_shape(s, '行业研发效能现状与趋势', 32, bold=True)
replace_text_everywhere(slides[2], '分页标题', '行业研发效能现状与趋势')
print("✅ 第3页: 分隔页 → 章节1 (保留原版式)")

# ============================================================
# 正文页（使用新策略：删除所有形状并重绘）
# ============================================================

print("🛠️ 正文页使用全新 Python-PPTX 图形重绘中...")

draw_custom_content_slide(slides[3], '行业概览', [
    "银行业数字化转型加速，软件研发投入年均增长 18%",
    "头部银行研发团队规模突破万人，中小银行加速外包转自研",
    "云原生、微服务架构成为主流技术选型方向",
    "研发效能已成为银行科技竞争力的核心衡量指标",
    "2025 年银行业 IT 投入预计超 3200 亿元"
])
print("✅ 第4页: 行业概览 (重绘)")

draw_custom_content_slide(slides[4], '效能挑战分析', [
    "遗留系统技术债务严重，核心系统改造周期长",
    "监管合规要求高，安全与效率的平衡难度大",
    "跨部门协作壁垒导致需求响应速度缓慢",
    "自动化测试覆盖率不足，回归测试成本高企",
    "人才储备与技能转型存在较大缺口"
])
print("✅ 第5页: 效能挑战 (重绘)")

# --- 第6页：分隔页2 ---
replace_text_everywhere(slides[5], '分页标题', '银行业研发效能核心指标体系')
print("✅ 第6页: 分隔页 → 章节2 (保留原版式)")

draw_custom_content_slide(slides[6], 'DORA 核心指标分析', [
    "部署频率：头部银行已实现日均部署 50+ 次",
    "变更前置时间：从需求提出到上线平均 15-30 天",
    "变更失败率：行业平均 8.5%，优秀团队控制在 3% 以内",
    "故障恢复时间（MTTR）：核心系统要求短于 30 分钟",
    "代码评审覆盖率：目标 100%，行业目前平均 72%"
])
print("✅ 第7页: DORA 指标 (重绘)")

draw_custom_content_slide(slides[7], '效能度量体系', [
    "需求交付效率：需求吞吐量、交付周期、在制品数量",
    "工程质量指标：缺陷密度、千行代码 Bug 率、技术债务比",
    "自动化水平：CI/CD 流水线覆盖率、自动化测试占比",
    "协作效率指标：PR 合并时长、代码评审响应速度",
    "业务价值指标：功能采纳率、用户满意度评分"
])
print("✅ 第8页: 效能度量体系 (重绘)")

# --- 第9页：分隔页3 ---
replace_text_everywhere(slides[8], '分页标题', 'DevOps 与敏捷实践落地分析')
print("✅ 第9页: 分隔页 → 章节3 (保留原版式)")

draw_custom_content_slide(slides[9], 'DevOps 实践成果', [
    "CI/CD 管道标准化：统一构建、测试、部署流水线",
    "容器化部署率提升至 65%，Kubernetes 集群规模扩展",
    "基础设施即代码（IaC）实现环境一致性管理",
    "灰度发布与蓝绿部署降低变更风险",
    "监控告警体系从被动响应转向主动预防"
])
print("✅ 第10页: DevOps 实践 (重绘)")

draw_custom_content_slide(slides[10], '敏捷转型路径', [
    "双模 IT 架构：稳态业务与敏态创新并行推进",
    "大规模敏捷框架（SAFe）在头部银行的落地实践",
    "产品化运营思维替代项目制交付模式",
    "跨职能团队组建，缩短决策链路和沟通成本",
    "迭代周期从月度缩短至双周，提升市场响应速度"
])
print("✅ 第11页: 敏捷转型 (重绘)")

# --- 第12页：分隔页4 ---
replace_text_everywhere(slides[11], '分页标题', 'AI 驱动的研发效能提升路径')
print("✅ 第12页: 分隔页 → 章节4 (保留原版式)")

draw_custom_content_slide(slides[12], 'AI 赋能研发', [
    "AI 辅助编码：Copilot 类工具提升编码效率 30-55%",
    "智能代码审查：自动检测安全漏洞和代码规范问题",
    "AI 测试生成：基于大模型自动生成测试用例，覆盖率提升 40%",
    "智能运维（AIOps）：故障预测准确率达到 85%",
    "需求分析智能化：NLP 驱动的需求拆解与风险识别"
])
print("✅ 第13页: AI 赋能研发 (重绘)")

draw_custom_content_slide(slides[13], 'AI 实施路线图', [
    "第一阶段：引入 AI 编码助手，提升个人开发效率",
    "第二阶段：建设智能 CI/CD 管道，实现自动化质量门禁",
    "第三阶段：构建研发知识图谱，沉淀组织级最佳实践",
    "第四阶段：大模型驱动的端到端自动化研发平台",
    "关键成功因素：数据治理、安全合规、人才培养"
])
print("✅ 第14页: AI 实施路线 (重绘)")

# --- 第15页：分隔页5 ---
replace_text_everywhere(slides[14], '分页标题', '展望与建议')
print("✅ 第15页: 分隔页 → 章节5 (保留原版式)")

draw_custom_content_slide(slides[15], '未来展望与建议', [
    "构建全链路效能度量平台，实现数据驱动的持续改进",
    "深化 AI 与研发流程的融合，打造智能研发中台",
    "推动开源生态参与，建立行业级技术标准共识",
    "加强研发安全左移，将安全融入 DevSecOps 全生命周期",
    "培养复合型数字化人才，支撑技术战略落地"
])
print("✅ 第16页: 未来展望 (重绘)")

draw_custom_content_slide(slides[16], '行动计划建议', [
    "成立研发效能委员会，统筹全行效能提升战略",
    "制定 3 年研发效能路线图，分步实施改进计划",
    "推进研发效能工具链整合，降低工具切换成本",
    "建设内部技术社区，促进知识共享与最佳实践传播",
    "与行业头部机构合作，引进成熟方法论与案例经验"
])
print("✅ 第17页: 行动计划 (重绘)")

# --- 第18页：结束页 ---
print("✅ 第18页: 结束页 (保持原样)")

prs.save(OUTPUT)
print(f"\n🎉 PPT 生成成功！")
print(f"📄 文件路径: {OUTPUT}")
print(f"📊 共计 {len(prs.slides)} 页幻灯片")

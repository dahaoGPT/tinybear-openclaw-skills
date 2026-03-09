"""
PPT 内容读取脚本
=================
读取 PPTX 文件并输出 JSON 格式的结构化信息。

用法:
    python read_ppt.py <file_path>
"""

import sys
import json
import os

# 确保可以导入 python-pptx
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("❌ python-pptx 未安装。请运行: pip install python-pptx")
    sys.exit(1)


def read_ppt(filepath):
    """读取 PPTX 文件并返回结构化数据"""
    if not os.path.exists(filepath):
        print(f"❌ 文件不存在: {filepath}")
        sys.exit(1)

    prs = Presentation(filepath)

    result = {
        "file": os.path.basename(filepath),
        "slide_count": len(prs.slides),
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i + 1,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "texts": [],
            "images": [],
            "charts": [],
            "tables": [],
        }

        for shape in slide.shapes:
            # 文本内容
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_info["texts"].append({
                        "name": shape.name,
                        "text": text,
                        "position": {
                            "left": round(shape.left / 914400, 2),
                            "top": round(shape.top / 914400, 2),
                        },
                        "size": {
                            "width": round(shape.width / 914400, 2),
                            "height": round(shape.height / 914400, 2),
                        },
                    })

            # 图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info["images"].append({
                    "name": shape.name,
                    "position": {
                        "left": round(shape.left / 914400, 2),
                        "top": round(shape.top / 914400, 2),
                    },
                    "size": {
                        "width": round(shape.width / 914400, 2),
                        "height": round(shape.height / 914400, 2),
                    },
                })

            # 图表
            if shape.has_chart:
                chart = shape.chart
                slide_info["charts"].append({
                    "name": shape.name,
                    "chart_type": str(chart.chart_type),
                    "has_legend": chart.has_legend,
                })

            # 表格
            if shape.has_table:
                table = shape.table
                rows_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    rows_data.append(row_data)
                slide_info["tables"].append({
                    "name": shape.name,
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "data": rows_data,
                })

        result["slides"].append(slide_info)

    return result


def main():
    if len(sys.argv) < 2:
        print("用法: python read_ppt.py <file_path>")
        sys.exit(1)

    filepath = sys.argv[1]
    result = read_ppt(filepath)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()

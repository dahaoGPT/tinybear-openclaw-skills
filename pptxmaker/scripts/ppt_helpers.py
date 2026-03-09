"""
PPT Helper Library - 简化 python-pptx 操作
==========================================
提供预设布局、配色方案和快捷方法，降低 PPT 生成的复杂度。

【重要规范】基于模板创建 PPT 时的注意事项：
1. 模板 slide 中的文本框绝大部分是手动添加的 shapes（非 layout placeholders），
   使用 add_slide(layout) 只会继承 layout 中的占位符，不会复制这些手动 shapes。
   必须使用 duplicate_slide() 深拷贝模板中的完整 slide，再替换文本内容。
2. 字体设置必须同时设置拉丁字体(font.name)和东亚字体(a:ea)：
   - 含中文字符的文本：font.name 设为中文字体（如"圆体-简"），a:ea 也设为中文字体
   - 纯英文/数字文本：font.name 设为英文字体（如"effra"），a:ea 设为中文字体（回退用）
   使用 set_run_font() 函数可自动处理。
"""

import re
import copy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn


# ============================================================
# 字体工具函数（模块级）
# ============================================================

def _has_cjk(text):
    """判断文本是否包含中日韩字符"""
    for ch in text:
        cp = ord(ch)
        if (0x4E00 <= cp <= 0x9FFF or    # CJK Unified Ideographs
            0x3400 <= cp <= 0x4DBF or    # CJK Extension A
            0x3000 <= cp <= 0x303F or    # CJK Symbols and Punctuation
            0xFF00 <= cp <= 0xFFEF or    # Fullwidth Forms
            0x2E80 <= cp <= 0x2EFF or    # CJK Radicals Supplement
            0xF900 <= cp <= 0xFAFF or    # CJK Compatibility Ideographs
            0xFE30 <= cp <= 0xFE4F):     # CJK Compatibility Forms
            return True
    return False


def set_run_font(run, font_cn='圆体-简', font_en='effra'):
    """
    为 run 智能设置字体：
    - 含中文字符时：font.name = font_cn, a:ea = font_cn
    - 纯英文/数字时：font.name = font_en, a:ea = font_cn（回退用）

    这确保 PowerPoint 和 WPS 都能正确渲染中英文字体。
    """
    text = run.text or ''
    if _has_cjk(text):
        run.font.name = font_cn
    else:
        run.font.name = font_en

    # 设置东亚字体（中文回退）
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_cn)


def duplicate_slide(prs, src_slide_index):
    """
    深拷贝模板中的指定 slide（包括所有手动添加的 shapes、图表、图片等）。

    这是基于模板创建 PPT 的核心方法，因为模板 slide 中的文本框、图表等
    绑定在 slide 本身而非 slide_layout 上，add_slide(layout) 无法复制它们。

    Args:
        prs: Presentation 对象
        src_slide_index: 源 slide 的索引号（0-based）

    Returns:
        新复制的 slide 对象
    """
    src_slide = prs.slides[src_slide_index]
    slide_layout = src_slide.slide_layout

    # 创建新 slide
    new_slide = prs.slides.add_slide(slide_layout)

    # 清除新 slide 中 layout 自动生成的占位符
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # 复制源 slide 的所有关系（图片、图表等），建立 rId 映射
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        reltype_suffix = rel.reltype.split('/')[-1]
        if reltype_suffix in ('slideLayout', 'notesSlide'):
            continue
        new_rId = new_slide.part.relate_to(rel._target, rel.reltype)
        rId_map[rId] = new_rId

    # 深拷贝源 slide 的所有 shape XML，并更新 rId 引用
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        for old_rId, new_rId in rId_map.items():
            if old_rId != new_rId:
                for attr_elem in el.iter():
                    for attr_name, attr_val in list(attr_elem.attrib.items()):
                        if attr_val == old_rId:
                            attr_elem.set(attr_name, new_rId)
                    r_id = attr_elem.get(qn('r:id'))
                    if r_id == old_rId:
                        attr_elem.set(qn('r:id'), new_rId)
                    r_link = attr_elem.get(qn('r:link'))
                    if r_link == old_rId:
                        attr_elem.set(qn('r:link'), new_rId)
                    r_embed = attr_elem.get(qn('r:embed'))
                    if r_embed == old_rId:
                        attr_elem.set(qn('r:embed'), new_rId)
        new_slide.shapes._spTree.append(el)

    return new_slide


def fix_content_title(shape, new_text=None, font_cn='圆体-简', font_en='effra'):
    """
    将正文页标题 shape 统一调整到母版规范位置：
    正文的标题需要放在左上角，在长方形图标之后，在横线之上，并把正文标题的字体设置成 24 级（24pt）居左展示。

    Args:
        shape: 标题 shape 对象
        new_text: 新标题文本（None 则不替换文本，只调整位置和字号）
        font_cn: 中文字体名
        font_en: 英文字体名
    """
    if shape is None:
        return
    shape.left = Inches(0.51)
    shape.top = Inches(0.13)
    shape.width = Inches(5.50)
    shape.height = Inches(0.49)
    if new_text is not None and shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            first_run = tf.paragraphs[0].runs[0]
            first_run.text = new_text
            set_run_font(first_run, font_cn, font_en)
            for run in tf.paragraphs[0].runs[1:]:
                run.text = ''
            for para in tf.paragraphs[1:]:
                for run in para.runs:
                    run.text = ''
        else:
            tf.paragraphs[0].text = new_text
            for run in tf.paragraphs[0].runs:
                set_run_font(run, font_cn, font_en)
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(24)


# ============================================================
# 配色方案
# ============================================================

THEMES = {
    "bank_brand": {
        "primary": RGBColor(255, 85, 5),      # 基础色: 橙红
        "secondary": RGBColor(244, 185, 0),    # 基础色: 黄色
        "accent": RGBColor(53, 4, 100),       # 基础色: 深紫 (原图为 53,4,64，微调以匹配更深紫色)
        "bg": RGBColor(255, 255, 255),        # 基础色: 白色
        "text": RGBColor(53, 4, 100),         # 正文颜色(深色)
        "light_bg": RGBColor(243, 157, 107),   # 辅助色: 浅橙
        "font_cn": "圆体-简",                  # 指定中文字体
        "font_en": "effra",                   # 指定英文字母/数字字体
    },
    "business_blue": {
        "primary": RGBColor(0x1B, 0x3A, 0x5C),
        "secondary": RGBColor(0x4A, 0x90, 0xD9),
        "accent": RGBColor(0xF3, 0x9C, 0x12),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "light_bg": RGBColor(0xEC, 0xF0, 0xF1),
        "font_cn": "Microsoft YaHei",
        "font_en": "Calibri",
    },
    "tech_dark": {
        "primary": RGBColor(0x0D, 0x1B, 0x2A),
        "secondary": RGBColor(0x1B, 0x99, 0x8B),
        "accent": RGBColor(0x00, 0xF5, 0xD4),
        "bg": RGBColor(0x0D, 0x1B, 0x2A),
        "text": RGBColor(0xE0, 0xE0, 0xE0),
        "light_bg": RGBColor(0x1B, 0x2A, 0x3A),
        "font_cn": "Source Han Sans CN",
        "font_en": "Consolas",
    },
    "fresh_green": {
        "primary": RGBColor(0x2D, 0x6A, 0x4F),
        "secondary": RGBColor(0x74, 0xC6, 0x9D),
        "accent": RGBColor(0x40, 0x91, 0x6C),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2D, 0x3E, 0x2D),
        "light_bg": RGBColor(0xD8, 0xF3, 0xDC),
        "font_cn": "Microsoft YaHei",
        "font_en": "Segoe UI",
    },
    "warm_orange": {
        "primary": RGBColor(0xE7, 0x6F, 0x51),
        "secondary": RGBColor(0xF4, 0xA2, 0x61),
        "accent": RGBColor(0x26, 0x46, 0x53),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x26, 0x46, 0x53),
        "light_bg": RGBColor(0xFD, 0xF0, 0xE2),
        "font_cn": "FZLanTingHei-R-GBK",
        "font_en": "Arial",
    },
    "elegant_purple": {
        "primary": RGBColor(0x43, 0x23, 0x71),
        "secondary": RGBColor(0x71, 0x4E, 0xBF),
        "accent": RGBColor(0xFA, 0xAE, 0x7B),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x33, 0x1A, 0x55),
        "light_bg": RGBColor(0xF0, 0xE6, 0xFA),
        "font_cn": "Source Han Serif CN",
        "font_en": "Georgia",
    },
    "minimal_gray": {
        "primary": RGBColor(0x2B, 0x2D, 0x42),
        "secondary": RGBColor(0x8D, 0x99, 0xAE),
        "accent": RGBColor(0xEF, 0x23, 0x3C),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2B, 0x2D, 0x42),
        "light_bg": RGBColor(0xED, 0xF2, 0xF4),
        "font_cn": "Microsoft YaHei",
        "font_en": "Helvetica",
    },
}


# ============================================================
# 布局坐标常量（单位：Inches）
# ============================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

LAYOUTS = {
    "title": {
        "title": {"left": Inches(1.0), "top": Inches(2.0), "width": Inches(11.333), "height": Inches(1.5)},
        "subtitle": {"left": Inches(2.0), "top": Inches(3.8), "width": Inches(9.333), "height": Inches(1.2)},
    },
    "content": {
        "title": {"left": Inches(0.7), "top": Inches(0.4), "width": Inches(11.933), "height": Inches(0.8)},
        "body": {"left": Inches(0.7), "top": Inches(1.5), "width": Inches(11.933), "height": Inches(5.5)},
    },
    "two_column": {
        "title": {"left": Inches(0.7), "top": Inches(0.4), "width": Inches(11.933), "height": Inches(0.8)},
        "left": {"left": Inches(0.7), "top": Inches(1.5), "width": Inches(5.767), "height": Inches(5.5)},
        "right": {"left": Inches(6.867), "top": Inches(1.5), "width": Inches(5.767), "height": Inches(5.5)},
    },
    "section": {
        "title": {"left": Inches(1.5), "top": Inches(2.5), "width": Inches(10.333), "height": Inches(1.5)},
        "subtitle": {"left": Inches(2.5), "top": Inches(4.2), "width": Inches(8.333), "height": Inches(0.8)},
    },
    "image": {
        "title": {"left": Inches(0.7), "top": Inches(0.4), "width": Inches(11.933), "height": Inches(0.8)},
        "image": {"left": Inches(1.5), "top": Inches(1.5), "width": Inches(10.333), "height": Inches(4.5)},
        "caption": {"left": Inches(1.5), "top": Inches(6.2), "width": Inches(10.333), "height": Inches(0.8)},
    },
    "chart": {
        "title": {"left": Inches(0.7), "top": Inches(0.4), "width": Inches(11.933), "height": Inches(0.8)},
        "chart": {"left": Inches(1.0), "top": Inches(1.5), "width": Inches(11.333), "height": Inches(5.5)},
    },
}


# ============================================================
# PPTHelper 类
# ============================================================

class PPTHelper:
    """PowerPoint 演示文稿生成辅助类"""

    def __init__(self, template_path=None, theme="bank_brand"):
        """
        初始化 PPT 辅助工具。

        Args:
            template_path: 可选，.pptx 模板文件路径
            theme: 配色方案名称，默认 'bank_brand'
        """
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT

        self.theme = THEMES.get(theme, THEMES["bank_brand"])
        self._slide_count = 0

    # ----------------------------------------------------------
    # 布局快捷方法
    # ----------------------------------------------------------

    def add_title_slide(self, title, subtitle=""):
        """添加标题页（封面页）"""
        slide = self._add_blank_slide()
        layout = LAYOUTS["title"]

        # 背景色块
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, SLIDE_HEIGHT, self.theme["primary"])

        # 标题
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=36, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

        # 副标题
        if subtitle:
            sub_box = self._add_textbox(slide, **layout["subtitle"])
            self._set_text(sub_box, subtitle, size=20,
                           color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

        return slide

    def add_content_slide(self, title, bullets):
        """
        添加内容页（标题 + 要点列表）。

        Args:
            title: 幻灯片标题
            bullets: 要点列表，例如 ["要点1", "要点2", ...]
        """
        slide = self._add_blank_slide()
        layout = LAYOUTS["content"]

        # 标题栏背景
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, Inches(1.3), self.theme["primary"])
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=28, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))

        # 要点内容
        body_box = self._add_textbox(slide, **layout["body"])
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.space_after = Pt(12)
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = self.theme["text"]
                set_run_font(run, self.theme["font_cn"], self.theme["font_en"])

        return slide

    def add_two_column_slide(self, title, left_title, left_items,
                              right_title, right_items):
        """
        添加两栏对比页。

        Args:
            title: 幻灯片标题
            left_title: 左栏标题
            left_items: 左栏要点列表
            right_title: 右栏标题
            right_items: 右栏要点列表
        """
        slide = self._add_blank_slide()
        layout = LAYOUTS["two_column"]

        # 标题栏
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, Inches(1.3), self.theme["primary"])
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=28, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))

        # 左栏
        self._add_column(slide, layout["left"], left_title, left_items)

        # 右栏
        self._add_column(slide, layout["right"], right_title, right_items)

        return slide

    def add_image_slide(self, title, image_path, caption=""):
        """
        添加图文页。

        Args:
            title: 幻灯片标题
            image_path: 图片文件路径
            caption: 图片说明文字
        """
        slide = self._add_blank_slide()
        layout = LAYOUTS["image"]

        # 标题
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, Inches(1.3), self.theme["primary"])
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=28, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))

        # 图片
        img_pos = layout["image"]
        slide.shapes.add_picture(image_path,
                                 img_pos["left"], img_pos["top"],
                                 img_pos["width"], img_pos["height"])

        # 说明文字
        if caption:
            cap_box = self._add_textbox(slide, **layout["caption"])
            self._set_text(cap_box, caption, size=14,
                           color=self.theme["secondary"], alignment=PP_ALIGN.CENTER)

        return slide

    def add_chart_slide(self, title, chart_type, categories, series_data):
        """
        添加图表页。

        Args:
            title: 幻灯片标题
            chart_type: 图表类型 ('bar', 'pie', 'line')
            categories: 分类标签列表，如 ['Q1', 'Q2', 'Q3', 'Q4']
            series_data: 数据系列字典，如 {'销售额': [100, 200, 300, 400]}
        """
        slide = self._add_blank_slide()
        layout = LAYOUTS["chart"]

        # 标题
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, Inches(1.3), self.theme["primary"])
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=28, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))

        # 图表
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data.items():
            chart_data.add_series(name, values)

        chart_types_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE_MARKERS,
        }
        xl_chart_type = chart_types_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_pos = layout["chart"]
        chart_frame = slide.shapes.add_chart(
            xl_chart_type, chart_pos["left"], chart_pos["top"],
            chart_pos["width"], chart_pos["height"], chart_data
        )

        # 简单样式
        chart = chart_frame.chart
        chart.has_legend = True if len(series_data) > 1 else False

        return slide

    def add_section_slide(self, title, subtitle=""):
        """添加章节过渡页"""
        slide = self._add_blank_slide()
        layout = LAYOUTS["section"]

        # 背景
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, SLIDE_HEIGHT, self.theme["secondary"])

        # 标题
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=36, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

        if subtitle:
            sub_box = self._add_textbox(slide, **layout["subtitle"])
            self._set_text(sub_box, subtitle, size=18,
                           color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

        return slide

    def add_summary_slide(self, title, key_points):
        """
        添加总结页。

        Args:
            title: 总结标题（如 "关键要点"）
            key_points: 要点列表
        """
        slide = self._add_blank_slide()
        layout = LAYOUTS["content"]

        # 标题栏
        self._add_shape_fill(slide, Inches(0), Inches(0),
                             SLIDE_WIDTH, Inches(1.3), self.theme["accent"])
        title_box = self._add_textbox(slide, **layout["title"])
        self._set_text(title_box, title, size=28, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))

        # 要点
        body_box = self._add_textbox(slide, **layout["body"])
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(key_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {point}"
            p.space_after = Pt(14)
            for run in p.runs:
                run.font.size = Pt(22)
                run.font.bold = True
                run.font.color.rgb = self.theme["primary"]
                set_run_font(run, self.theme["font_cn"], self.theme["font_en"])

        return slide

    # ----------------------------------------------------------
    # 样式方法
    # ----------------------------------------------------------

    def set_font(self, run, size=18, bold=False, color=None, font_name=None):
        """
        统一设置字体属性（自动处理中英文双字体）。

        Args:
            run: pptx Run 对象
            size: 字号（pt）
            bold: 是否加粗
            color: RGBColor 对象
            font_name: 指定英文字体名称（中文字体从 theme 中获取）
        """
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        set_run_font(run, self.theme["font_cn"], font_name or self.theme["font_en"])

    def apply_theme(self, theme_name):
        """切换配色方案"""
        if theme_name in THEMES:
            self.theme = THEMES[theme_name]

    # ----------------------------------------------------------
    # 工具方法
    # ----------------------------------------------------------

    def save(self, filename):
        """保存文件，自动添加 .pptx 后缀"""
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        self.prs.save(filename)
        return filename

    @staticmethod
    def hex_color(hex_str):
        """将十六进制颜色字符串转为 RGBColor。
        示例: hex_color('#1B3A5C') → RGBColor(0x1B, 0x3A, 0x5C)
        """
        hex_str = hex_str.lstrip("#")
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)

    # ----------------------------------------------------------
    # 内部辅助方法
    # ----------------------------------------------------------

    def _add_blank_slide(self):
        """添加空白幻灯片"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)
        self._slide_count += 1
        return slide

    def _add_textbox(self, slide, left, top, width, height):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        return txBox

    def _set_text(self, textbox, text, size=18, bold=False, color=None,
                  alignment=PP_ALIGN.LEFT, font_name=None):
        """设置文本框内容和样式"""
        tf = textbox.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].alignment = alignment

        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            set_run_font(run, self.theme["font_cn"], font_name or self.theme["font_en"])

    def _add_shape_fill(self, slide, left, top, width, height, color):
        """添加纯色填充矩形（用作背景色块）"""
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()  # 无边框
        return shape

    def _add_column(self, slide, pos, col_title, items):
        """添加一列内容（用于两栏布局）"""
        # 列标题
        title_height = Inches(0.6)
        title_box = self._add_textbox(slide, pos["left"], pos["top"],
                                      pos["width"], title_height)
        self._set_text(title_box, col_title, size=22, bold=True,
                       color=self.theme["primary"])

        # 列内容
        body_top = pos["top"] + title_height + Inches(0.2)
        body_height = pos["height"] - title_height - Inches(0.2)
        body_box = self._add_textbox(slide, pos["left"], body_top,
                                     pos["width"], body_height)
        tf = body_box.text_frame
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = self.theme["text"]
                set_run_font(run, self.theme["font_cn"], self.theme["font_en"])


# ============================================================
# 便捷入口
# ============================================================

def quick_create(title, slides_config, theme="bank_brand",
                 template=None, output="presentation.pptx"):
    """
    快速创建 PPT 的便捷函数。

    Args:
        title: 演示文稿标题
        slides_config: 幻灯片配置列表，每项是 dict，含 'type' 和对应参数
        theme: 配色方案名称
        template: 可选，模板文件路径
        output: 输出文件名

    示例:
        quick_create("我的演示", [
            {"type": "title", "title": "欢迎", "subtitle": "副标题"},
            {"type": "content", "title": "要点", "bullets": ["A", "B", "C"]},
            {"type": "summary", "title": "总结", "key_points": ["重点1", "重点2"]},
        ])
    """
    helper = PPTHelper(template_path=template, theme=theme)

    # 封面
    helper.add_title_slide(title)

    # 各页
    for cfg in slides_config:
        slide_type = cfg.get("type", "content")

        if slide_type == "title":
            helper.add_title_slide(cfg["title"], cfg.get("subtitle", ""))
        elif slide_type == "content":
            helper.add_content_slide(cfg["title"], cfg.get("bullets", []))
        elif slide_type == "two_column":
            helper.add_two_column_slide(
                cfg["title"],
                cfg.get("left_title", ""), cfg.get("left_items", []),
                cfg.get("right_title", ""), cfg.get("right_items", [])
            )
        elif slide_type == "section":
            helper.add_section_slide(cfg["title"], cfg.get("subtitle", ""))
        elif slide_type == "image":
            helper.add_image_slide(cfg["title"], cfg["image_path"],
                                   cfg.get("caption", ""))
        elif slide_type == "chart":
            helper.add_chart_slide(cfg["title"], cfg.get("chart_type", "bar"),
                                   cfg.get("categories", []),
                                   cfg.get("series_data", {}))
        elif slide_type == "summary":
            helper.add_summary_slide(cfg["title"], cfg.get("key_points", []))

    return helper.save(output)

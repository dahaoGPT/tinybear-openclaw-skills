"""
模拟 /ppt-maker:from-template company-2026 银行业软件研发效能报告
先读取模板结构，输出可用的布局信息
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

TEMPLATE = r'e:\2-AI\6-workspaces\managertools\claude-code\pptxmaker\templates\company-2026.pptx'

print("=== 读取模板结构 ===")
prs = Presentation(TEMPLATE)
print(f"幻灯片尺寸: {prs.slide_width} x {prs.slide_height} EMU")
print(f"  即: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} 英寸")
print(f"现有幻灯片数量: {len(prs.slides)}")

print("\n=== 可用母版布局 ===")
for i, layout in enumerate(prs.slide_layouts):
    ph_info = []
    for ph in layout.placeholders:
        ph_info.append(f"    占位符[{ph.placeholder_format.idx}]: {ph.name} (类型:{ph.placeholder_format.type})")
    print(f"布局[{i}]: {layout.name}")
    for p in ph_info:
        print(p)

print("\n=== 现有幻灯片内容摘要 ===")
for idx, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else "未知"
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t[:50])
    print(f"第{idx+1}页 [布局: {layout_name}]: {'; '.join(texts[:3]) if texts else '(无文本)'}")

print("\n=== 模板读取完成 ===")
